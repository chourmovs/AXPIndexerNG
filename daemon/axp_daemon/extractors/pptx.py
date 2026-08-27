def extract(path):
    from pptx import Presentation
    return [('\n'.join(s.text for shape in slide.shapes if hasattr(shape,'text_frame') for s in shape.text_frame.paragraphs),i+1) for i,slide in enumerate(Presentation(path).slides)]
